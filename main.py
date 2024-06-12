import sys
from PyQt5.QtWidgets import (QApplication, QMainWindow, QFileDialog, QTextEdit, QVBoxLayout, 
                             QHBoxLayout, QWidget, QPushButton, QLineEdit, QLabel, QMessageBox)
from PyQt5.QtGui import QFont, QIcon
from PyQt5.QtCore import Qt
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
import google.generativeai as genai
from datetime import datetime

# Configurar la API KEY
GOOGLE_API_KEY = 'AIzaSyAMwbG-N2bpQ6pzNo54oY0FWOgsC6xlizY'
genai.configure(api_key=GOOGLE_API_KEY)

# Función para rebajar el tamaño de la respuesta
def rebajar(text):
    text = text.replace('•', ' *')
    return text

# Inicializar el modelo generativo de la IA
modelo = genai.GenerativeModel('gemini-pro')

class App(QMainWindow):
    def __init__(self):
        super().__init__()
        self.title = 'Consultas de Archivos de Investigación'
        self.left = 100
        self.top = 100
        self.width = 800
        self.height = 600
        self.texto_documento = ""
        self.historial = []
        self.historial_guardado = False
        self.preguntas_realizadas = 0
        self.documento_cargado = False  # Variable para verificar si se ha cargado un documento

        self.initUI()

    def initUI(self):
        self.setWindowTitle(self.title)
        self.setGeometry(self.left, self.top, self.width, self.height)

        # Establecer icono de la aplicación
        self.setWindowIcon(QIcon('Logo_Uneg.png'))
        
        central_widget = QWidget(self)
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout()

        # Botón de cargar archivo
        self.btn_cargar = QPushButton('Cargar Archivo', self)
        self.btn_cargar.clicked.connect(self.cargar_archivo)
        self.btn_cargar.setFixedWidth(120)
        main_layout.addWidget(self.btn_cargar)

        # Layout horizontal para la pregunta y el botón de preguntar
        pregunta_layout = QHBoxLayout()
        
        self.lbl_pregunta = QLabel('Pregunta:', self)
        pregunta_layout.addWidget(self.lbl_pregunta)
        
        self.entrada_pregunta = QLineEdit(self)
        pregunta_layout.addWidget(self.entrada_pregunta)
        
        self.btn_preguntar = QPushButton('Preguntar', self)
        self.btn_preguntar.clicked.connect(self.hacer_pregunta)
        self.btn_preguntar.setFixedWidth(120)
        pregunta_layout.addWidget(self.btn_preguntar)

        main_layout.addLayout(pregunta_layout)

        # Área de texto para el historial
        self.texto_historial = QTextEdit(self)
        self.texto_historial.setReadOnly(True)
        self.texto_historial.setFont(QFont("Times New Roman", 12))
        main_layout.addWidget(self.texto_historial)

        # Layout horizontal para los botones de guardar historial y salir
        botones_layout = QHBoxLayout()
        botones_layout.addStretch()  # Añadir un espacio flexible
        
        self.btn_guardar = QPushButton('Guardar Historial', self)
        self.btn_guardar.clicked.connect(self.guardar_historial)
        self.btn_guardar.setFixedWidth(120)
        botones_layout.addWidget(self.btn_guardar)
        
        self.btn_salir = QPushButton('Salir', self)
        self.btn_salir.clicked.connect(self.salir)
        self.btn_salir.setFixedWidth(120)
        botones_layout.addWidget(self.btn_salir)

        main_layout.addLayout(botones_layout)
        
        central_widget.setLayout(main_layout)

        # Estilos de modo oscuro claro suave
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f0f0f0;
            }
            QPushButton {
                background-color: #cccccc;
                color: #000000;
                border: 1px solid #888888;
                border-radius: 5px;
                padding: 5px;
            }
            QPushButton:hover {
                background-color: #bbbbbb;
            }
            QLabel {
                color: #000000;
            }
            QLineEdit {
                background-color: #ffffff;
                color: #000000;
                border: 1px solid #888888;
                padding: 5px;
            }
            QTextEdit {
                background-color: #ffffff;
                color: #000000;
                border: 1px solid #888888;
                padding: 5px;
            }
        """)

    def cargar_archivo(self):
        options = QFileDialog.Options()
        file, _ = QFileDialog.getOpenFileName(self, "Cargar Archivo", "", "Todos los archivos (*.docx *.pdf *.pptx);;Documentos DOCX (*.docx);;Archivos PDF (*.pdf);;Presentaciones PowerPoint (*.pptx)", options=options)
        if file:
            extension = file.split('.')[-1].lower()
            if extension == "docx":
                doc = Document(file)
                self.texto_documento = ""
                for paragraph in doc.paragraphs:
                    self.texto_documento += paragraph.text + "\n"
            elif extension == "pdf":
                self.texto_documento = ""
                pdf_documento = fitz.open(file)
                for page_num in range(len(pdf_documento)):
                    page = pdf_documento.load_page(page_num)
                    self.texto_documento += page.get_text("text") + "\n"
            elif extension == "pptx":
                self.texto_documento = ""
                ppt_documento = Presentation(file)
                for slide in ppt_documento.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            self.texto_documento += shape.text + "\n"
            self.documento_cargado = True  # Marcar que se ha cargado un documento
            QMessageBox.information(self, "Archivo cargado", "El archivo se ha cargado correctamente.")
        else:
            QMessageBox.warning(self, "Carga fallida", "No se ha seleccionado ningún archivo.")

    def hacer_pregunta(self):
        if not self.documento_cargado:
            QMessageBox.warning(self, "No hay documento cargado", "Por favor, cargue un documento antes de hacer una pregunta.")
            return
        
        pregunta = self.entrada_pregunta.text()
        if pregunta.strip():
            self.preguntas_realizadas += 1
            pregunta_con_contexto = f"{self.texto_documento}\n{pregunta}"
            respuesta = modelo.generate_content(pregunta_con_contexto)
            respuesta_texto = respuesta.text
            self.historial.append(f"Pregunta: {pregunta}\nRespuesta: {respuesta_texto}\n")
            self.texto_historial.append(f"Pregunta: {pregunta}")
            self.texto_historial.append(f"Respuesta: {rebajar(respuesta_texto)}\n")
            self.entrada_pregunta.clear()
        else:
            QMessageBox.warning(self, "Pregunta vacía", "Por favor, ingrese una pregunta.")

    def guardar_historial(self):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        nombre_archivo = f"historial_conversacion_{timestamp}.txt"
        with open(nombre_archivo, 'w') as file:
            file.writelines(self.historial)
        QMessageBox.information(self, "Historial guardado", f"Historial guardado en '{nombre_archivo}'")
        self.historial_guardado = True

    def salir(self):
        if self.preguntas_realizadas > 0 and not self.historial_guardado:
            reply = QMessageBox.question(self, 'Advertencia', "Estás cerrando el programa sin haber guardado. ¿Estás seguro?", QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
            if reply == QMessageBox.Yes:
                self.close()
        else:
            self.close()

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = App()
    ex.show()
    sys.exit(app.exec_())
